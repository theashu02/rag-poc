import os
import time
import hashlib
from typing import List, Dict
from dataclasses import dataclass
from concurrent.futures import ThreadPoolExecutor
from tenacity import retry, stop_after_attempt, wait_exponential
from pptx import Presentation
from docx import Document
import textract

import spacy
import yake
import tiktoken
from keybert import KeyBERT
from langchain_text_splitters import RecursiveCharacterTextSplitter
from openai import OpenAI
from pinecone import Pinecone, ServerlessSpec
import pdfplumber
import fitz
import pytesseract
from PIL import Image
import io


# --- Configuration ---
@dataclass
class RAGConfig:
    openai_api_key: str = os.getenv("OPENAI_API_KEY")
    pinecone_api_key: str = os.getenv("PINECONE_API_KEY")
    index_name: str = os.getenv("PINECONE_INDEX")
    embedding_model: str = os.getenv("OPENAI_EMBEDDING_MODEL", "text-embedding-3-large")
    keyword_model: str = "all-MiniLM-L6-v2"
    chunk_size: int = 500
    chunk_overlap: int = 100
    batch_size: int = int(os.getenv("BATCH_SIZE", "50"))
    cache_dir: str = "/tmp/cache"
    max_workers: int = int(os.getenv("MAX_WORKERS", "6"))
    # Add timeout settings
    openai_timeout: int = int(os.getenv("OPENAI_TIMEOUT", "30"))
    pinecone_timeout: int = int(os.getenv("PINECONE_TIMEOUT", "30"))

config = RAGConfig()

# --- Initialization with lazy loading ---
_client = None
_pc = None
_nlp = None
_kw_extractor = None
_keybert_model = None
_encoding = None

def get_openai_client():
    global _client
    if _client is None:
        _client = OpenAI(api_key=config.openai_api_key, timeout=config.openai_timeout)
    return _client

def get_pinecone_index():
    global _pc
    if _pc is None:
        _pc = Pinecone(api_key=config.pinecone_api_key)
        
    dimension = 1536  # for text-embedding-3-small
    if config.index_name not in _pc.list_indexes().names():
        print(f"Creating index {config.index_name}...")
        _pc.create_index(
            name=config.index_name,
            dimension=dimension,
            metric="cosine",
            spec=ServerlessSpec(cloud="aws", region="us-east-1")
        )
        while not _pc.describe_index(config.index_name).status['ready']:
            time.sleep(1)
    return _pc.Index(config.index_name)

def get_spacy_nlp():
    global _nlp
    if _nlp is None:
        _nlp = spacy.load("en_core_web_lg")
    return _nlp

def get_yake_extractor():
    global _kw_extractor
    if _kw_extractor is None:
        _kw_extractor = yake.KeywordExtractor(n=3, top=15, dedupLim=0.9)
    return _kw_extractor

def get_keybert_model():
    global _keybert_model
    if _keybert_model is None:
        _keybert_model = KeyBERT(model=config.keyword_model)
    return _keybert_model

def get_tiktoken_encoding():
    global _encoding
    if _encoding is None:
        _encoding = tiktoken.encoding_for_model("gpt-4o")
    return _encoding

# Create cache directory
os.makedirs(config.cache_dir, exist_ok=True)

# --- Document Readers ---
def read_pdf(path):
    text_by_page: Dict[int, str] = {}
    images_for_ocr = []
    try:
        doc = fitz.open(path)
        for idx, page in enumerate(doc):
            page_num = idx + 1
            txt = page.get_text().strip()
            text_by_page[page_num] = txt

            if len(txt) < 50:
                for img in page.get_images():
                    pix = fitz.Pixmap(doc, img[0])
                    if pix.n - pix.alpha < 4:
                        images_for_ocr.append(
                            {"page": page_num, "data": pix.tobytes("png")}
                        )
        doc.close()
    except Exception as exc:                          
        print(f"PyMuPDF extraction failed: {exc}")

    for img in images_for_ocr:
        try:
            pil_img = Image.open(io.BytesIO(img["data"]))
            ocr_txt = pytesseract.image_to_string(
                pil_img, config="--oem 3 --psm 6"
            ).strip()
            if ocr_txt:
                merged = f"{text_by_page.get(img['page'], '')}\n{ocr_txt}".strip()
                text_by_page[img["page"]] = merged
        except Exception as exc:                    
            print(f"OCR failed (page {img['page']}): {exc}")

    try:
        with pdfplumber.open(path) as pdf:
            for idx, page in enumerate(pdf.pages):
                page_num = idx + 1
                extra = (page.extract_text() or "").strip()
                if extra and extra not in text_by_page.get(page_num, ""):
                    text_by_page[page_num] = (
                        f"{text_by_page.get(page_num, '')}\n{extra}".strip()
                    )
    except Exception as exc:
        print(f"pdfplumber extraction failed: {exc}")

    if text_by_page:
        full = "\n".join(
            text_by_page[p] for p in sorted(text_by_page) if text_by_page[p]
        )
        return normalize_text(full) if full else None
    return None

def read_txt(path):
    try:
        with open(path, "r", encoding='utf-8', errors="ignore") as f:
            return normalize_text(f.read())
    except Exception:
        return None

def read_json(path):
    try:
        with open(path, 'r', encoding='utf-8', errors="ignore") as f:
            return normalize_text(f.read())
    except Exception:
        return None

def read_pptx(path):
    try:
        prs = Presentation(path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)
    except Exception:
        return None

def read_docx(path):
    try:
        doc = Document(path)
        return "\n".join([para.text for para in doc.paragraphs])
    except Exception:
        return None

def read_doc(path):
    try:
        return textract.process(path).decode('utf-8', errors='ignore')
    except Exception:
        return None

def normalize_text(s: str) -> str:
    s = s.replace("\r\n", "\n").replace("\r", "\n")
    s = s.replace('\x00', '').replace('\xa0', ' ')
    s = s.encode('utf-8', 'ignore').decode('utf-8')
    return "\n".join([line.strip() for line in s.splitlines() if line.strip()])

# --- Text Processing with parallelization ---
KEEP_LABELS = {
    "PERSON", "ORG", "GPE", "PRODUCT", "EVENT", "NORP", "FAC","LOC","WORK_OF_ART", "LAW", "LANGUAGE", "DATE", "TIME", "MONEY", "PERCENT",
}

class EnhancedTextProcessor:
    def extract_entities(self, text: str) -> List[str]:
        nlp = get_spacy_nlp()
        doc = nlp(text[:1000000])
        entities = [ent.text for ent in doc.ents if ent.label_ in KEEP_LABELS]
        return list(set(entities))

    def extract_keywords(self, text: str) -> List[str]:
        keywords = set()
        
        # Extract YAKE keywords
        yake_kws = [kw for kw, _ in get_yake_extractor().extract_keywords(text)]
        keywords.update(yake_kws[:5])
        
        # Extract KeyBERT keywords
        try:
            keybert_kws = get_keybert_model().extract_keywords(
                text, keyphrase_ngram_range=(1, 3), stop_words='english', top_n=5
            )
            keywords.update([kw for kw, _ in keybert_kws])
        except Exception:
            pass
            
        return list(keywords)[:10]

# Initialize processor
processor = EnhancedTextProcessor()

# --- Chunking ---
class SemanticChunker:
    def __init__(self):
        encoding = get_tiktoken_encoding()
        self.splitter = RecursiveCharacterTextSplitter(
            chunk_size=config.chunk_size,
            chunk_overlap=config.chunk_overlap,
            length_function=lambda text: len(encoding.encode(text)),
            separators=["\n\n", "\n", ". ", " ", ""]
        )

    def chunk_with_context(self, text: str, metadata: Dict) -> List[Dict]:
        chunks = self.splitter.split_text(text)
        enhanced_chunks = []
        for idx, chunk_text in enumerate(chunks):
            enhanced_chunks.append({
                "text": chunk_text,
                "chunk_index": idx,
                "total_chunks": len(chunks),
                **metadata
            })
        return enhanced_chunks

chunker = SemanticChunker()

# --- Embedding Generation with retries ---
class EmbeddingGenerator:
    @retry(stop=stop_after_attempt(3), wait=wait_exponential(multiplier=1, min=4, max=10))
    def get_embeddings_batch(self, texts: List[str]) -> List[List[float]]:
        if not texts:
            return []
        try:
            client = get_openai_client()
            response = client.embeddings.create(
                model=config.embedding_model,
                input=texts,
                encoding_format="float"
            )
            return [data.embedding for data in response.data]
        except Exception as e:
            print(f"Batch embedding error: {e}")
            raise

embedder = EmbeddingGenerator()

# --- Main Processing Pipeline with parallelization ---
def process_file(file_path: str, original_filename: str) -> List[Dict]:
    print(f"Processing: {original_filename}")
    ext = os.path.splitext(original_filename)[1].lower()
    text = None
    if ext == ".pdf":
        text = read_pdf(file_path)
    elif ext == ".txt":
        text = read_txt(file_path)

    if not text:
        return []

    base_metadata = {
        "source": original_filename,
        "file_type": ext[1:],
        "processing_timestamp": time.time()
    }
    
    chunks = chunker.chunk_with_context(text, base_metadata)
    processed_chunks = []
    
    # Process chunks in parallel
    with ThreadPoolExecutor(max_workers=config.max_workers) as executor:
        futures = []
        for chunk_data in chunks:
            if len(chunk_data["text"]) < 50:
                continue
            
            # Submit chunk processing to thread pool
            future = executor.submit(process_single_chunk, chunk_data, original_filename)
            futures.append(future)
        
        # Collect results
        for future in futures:
            try:
                processed_chunks.append(future.result())
            except Exception as e:
                print(f"Error processing chunk: {e}")
    
    return processed_chunks

def process_single_chunk(chunk_data, original_filename):
    """Process a single chunk - extracted for parallel execution"""
    chunk_hash = hashlib.sha256(chunk_data["text"].encode()).hexdigest()[:16]
    chunk_id = f"{original_filename}-{chunk_hash}"
    
    metadata = {
        **chunk_data,
        "chunk_keywords": processor.extract_keywords(chunk_data["text"])
    }
    
    return {
        "id": chunk_id,
        "text": chunk_data["text"],
        "metadata": metadata
    }

# --- Pinecone Upload with batching and retries ---
@retry(stop=stop_after_attempt(3), wait=wait_exponential(multiplier=1, min=4, max=10))
def upload_batch_to_pinecone(batch, namespace):
    """Upload a single batch to Pinecone with retry logic"""
    index = get_pinecone_index()
    index.upsert(vectors=batch, namespace=namespace)
    return len(batch)

def upload_to_pinecone(chunks: List[Dict], namespace: str) -> int:
    if not chunks:
        return 0

    texts = [c["text"] for c in chunks]
    embeddings = embedder.get_embeddings_batch(texts)
    
    valid_chunks = []
    for chunk, embedding in zip(chunks, embeddings):
        if embedding:
            metadata = chunk["metadata"]
            metadata['text'] = chunk["text"]
            valid_chunks.append({
                "id": chunk["id"],
                "values": embedding,
                "metadata": metadata
            })

    if not valid_chunks:
        return 0
    
    uploaded_count = 0
    # Process batches in parallel
    with ThreadPoolExecutor(max_workers=config.max_workers) as executor:
        futures = []
        for i in range(0, len(valid_chunks), config.batch_size):
            batch = valid_chunks[i:i+config.batch_size]
            future = executor.submit(upload_batch_to_pinecone, batch, namespace)
            futures.append(future)
        
        # Collect results
        for future in futures:
            try:
                uploaded_count += future.result()
                print(f"Uploaded batch to namespace '{namespace}'.")
            except Exception as e:
                print(f"Pinecone upload error: {e}")
    
    return uploaded_count