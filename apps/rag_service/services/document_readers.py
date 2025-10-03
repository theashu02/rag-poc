from __future__ import annotations

import io
from typing import Dict, Optional

import fitz
import pdfplumber
import pytesseract
from PIL import Image
from pptx import Presentation
from docx import Document

from .text_normalizer import normalize_json, format_blocks


def _normalize_text(s: str) -> str:
    s = s.replace("\r\n", "\n").replace("\r", "\n")
    s = s.replace("\x00", "").replace("\xa0", " ")
    s = s.encode("utf-8", "ignore").decode("utf-8")
    return "\n".join([line.strip() for line in s.splitlines() if line.strip()])


def read_pdf(path: str) -> Optional[str]:
    text_by_page: Dict[int, str] = {}
    images_for_ocr = []

    try:
        with fitz.open(path) as doc:
            for idx, page in enumerate(doc):
                page_num = idx + 1
                txt = (page.get_text() or "").strip()
                text_by_page[page_num] = txt

                if len(txt) < 50:
                    for img in page.get_images(full=True):
                        pix = None
                        try:
                            pix = fitz.Pixmap(doc, img[0])
                            if pix.n - pix.alpha < 4:
                                images_for_ocr.append({
                                    "page": page_num,
                                    "data": pix.tobytes("png")
                                })
                        except Exception as exc:
                            print(f"[PDF] image extraction failed (page {page_num}): {exc}")
                        finally:
                            if hasattr(pix, "close"):
                                pix.close()
    except Exception as exc:
        print(f"[PDF] PyMuPDF extraction failed: {exc}")

    for img in images_for_ocr:
        try:
            pil_img = Image.open(io.BytesIO(img["data"]))
            ocr_txt = pytesseract.image_to_string(
                pil_img,
                config="--oem 3 --psm 6",
            ).strip()
            if ocr_txt:
                merged = f"{text_by_page.get(img['page'], '')}\n{ocr_txt}".strip()
                text_by_page[img["page"]] = merged
        except Exception as exc:
            print(f"[PDF] OCR failed (page {img['page']}): {exc}")

    try:
        with pdfplumber.open(path) as pdf:
            for idx, page in enumerate(pdf.pages):
                page_num = idx + 1
                extra = (page.extract_text() or "").strip()
                if extra and extra not in text_by_page.get(page_num, ""):
                    existing = text_by_page.get(page_num, "")
                    text_by_page[page_num] = f"{existing}\n{extra}".strip()
    except Exception as exc:
        print(f"[PDF] pdfplumber extraction failed: {exc}")

    if not text_by_page:
        return None

    full_text = "\n".join(
        text_by_page[p] for p in sorted(text_by_page) if text_by_page[p]
    )
    if not full_text:
        return None

    return _normalize_text(full_text)


def read_txt(path: str) -> Optional[str]:
    try:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            return _normalize_text(f.read())
    except Exception as exc:
        print(f"[TXT] read failed: {exc}")
        return None


def read_json(path: str) -> Optional[str]:
    try:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            cleaned_json = normalize_json(f.read())
            return format_blocks(cleaned_json, show_list_index=False)
    except Exception as exc:
        print(f"[JSON] read failed: {exc}")
        return None


def read_pptx(path: str) -> Optional[str]:
    try:
        prs = Presentation(path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text.append(shape.text)
        return "\n".join(text)
    except Exception as exc:
        print(f"[PPTX] read failed: {exc}")
        return None


def read_docx(path: str) -> Optional[str]:
    try:
        doc = Document(path)
        paragraphs = [para.text for para in doc.paragraphs if para.text]
        return "\n".join(paragraphs)
    except Exception as exc:
        print(f"[DOCX] read failed: {exc}")
        return None
